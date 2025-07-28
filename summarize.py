import os
import fitz
import subprocess
from pptx import Presentation
import openpyxl

def read_file_content(filepath):
    ext = os.path.splitext(filepath)[1].lower()

    if ext == '.txt':
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
        
    elif ext =='.pdf':
        doc = fitz.open(filepath)
        text = ""

        for page in doc:
            text += page.get_text()
        doc.close()
        return text
    
    elif ext in ['.java', '.c', '.py', '.cpp', '.js', '.html', '.css', '.json', '.xml']:
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(filepath, 'r', encoding='latin-1') as f:
                return f.read()

    elif ext in ['.ppt', '.pptx']:
        prs = Presentation(filepath)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"

        return text

    elif ext in ['.xlsx', '.xls']:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        text = ""

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text += f"\n--- Sheet: {sheet_name} ---\n"

            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_data):
                    text += " | ".join(row_data) + "\n"
        wb.close()
        return text            
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    
def chunk_text(text, chunk_size=500):
    words = text.split()

    return [" ".join(words[i : i+chunk_size]) for i in range(0, len(words), chunk_size)]

def summarize_text_from_file(filepath):
    content = read_file_content(filepath)
    chunks = chunk_text(content)

    all_summaries = []

    for chunk in chunks:
        ext = os.path.splitext(filepath)[1].lower()
        if ext in ['.py', '.java', '.cpp', '.c', '.js', '.html', '.css']:
            prompt = f"Summarize and explain this code:\n\n{chunk}"
        else:
            prompt = f"Summarize the text:\n\n{chunk}"

        prompt = f"Summarize the text: \n\n{chunk}"
        result = subprocess.run(
            ['ollama', 'run', 'mistral'],
            input = prompt.encode(),
            stdout = subprocess.PIPE,
            stderr = subprocess.PIPE
        )
        response = result.stdout.decode().strip()
        all_summaries.append(response)

    return "\n".join(all_summaries)    